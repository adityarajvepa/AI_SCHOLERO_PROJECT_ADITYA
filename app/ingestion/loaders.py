"""PDF and PPTX text extraction."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation


@dataclass(frozen=True)
class TextUnit:
    """One page or slide of extracted text."""

    unit_type: str  # "page" | "slide"
    unit_number: int  # 1-based
    text: str


def load_pdf_units(path: Path) -> list[TextUnit]:
    doc = fitz.open(path)
    units: list[TextUnit] = []
    try:
        for i in range(doc.page_count):
            page = doc.load_page(i)
            text = page.get_text("text") or ""
            units.append(TextUnit(unit_type="page", unit_number=i + 1, text=text))
    finally:
        doc.close()
    return units


def _shape_text(shape) -> str:  # type: ignore[no-untyped-def]
    parts: list[str] = []
    if hasattr(shape, "text") and shape.text:
        parts.append(shape.text)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:  # type: ignore[attr-defined]
            parts.append(_shape_text(child))
    return "\n".join(p for p in parts if p)


def load_pptx_units(path: Path) -> list[TextUnit]:
    prs = Presentation(str(path))
    units: list[TextUnit] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            texts.append(_shape_text(shape))
        raw = "\n".join(t for t in texts if t.strip())
        units.append(TextUnit(unit_type="slide", unit_number=idx, text=raw))
    return units
